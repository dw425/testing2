#!/usr/bin/env python3
"""Generate LHO Lite overview PowerPoint deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Blueprint color palette
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
BLUE = RGBColor(0x4B, 0x7B, 0xF5)
GREEN = RGBColor(0x34, 0xD3, 0x99)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
CYAN = RGBColor(0x00, 0xE5, 0xFF)
WHITE = RGBColor(0xE6, 0xED, 0xF3)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
BORDER = RGBColor(0x27, 0x2D, 0x3F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, size=16, bold=False, color=WHITE, spacing_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_before = Pt(spacing_before)
    return p


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_bullet_card(slide, left, top, width, height, title, items, title_color=BLUE):
    add_card(slide, left, top, width, height)
    tf = add_text(slide, left + 0.2, top + 0.15, width - 0.4, 0.4, title, size=14, bold=True, color=title_color)
    for item in items:
        add_para(tf, item, size=12, color=WHITE, spacing_before=4)


def footer(slide):
    add_text(slide, 0.5, 7.0, 5, 0.3, "LHO Lite by Blueprint Technologies", size=10, color=GRAY)
    add_text(slide, 8, 7.0, 5, 0.3, "Confidential", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1: Title
# =====================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s1)

# Blue accent bar at top
bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_text(s1, 1, 1.8, 11, 1.0, "LHO Lite", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, 1, 2.8, 11, 0.6, "Lakehouse Optimizer Lite", size=28, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s1, 1, 3.5, 11, 0.5, "Comprehensive Databricks Workspace Analysis", size=20, color=GRAY, align=PP_ALIGN.CENTER)

# Divider line
div = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.3), Inches(5.333), Inches(0.02))
div.fill.solid()
div.fill.fore_color.rgb = BORDER
div.line.fill.background()

# Key stats row
stats = [
    ("14", "Dashboard Tabs"),
    ("6", "System Tables"),
    ("22+", "API Endpoints"),
    ("4", "Compliance Frameworks"),
    ("16", "Security Findings"),
]
for i, (val, label) in enumerate(stats):
    x = 1.2 + i * 2.3
    add_text(s1, x, 4.6, 2, 0.5, val, size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s1, x, 5.1, 2, 0.3, label, size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s1, 1, 6.0, 11, 0.4, "by Blueprint Technologies", size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s1, 1, 6.4, 11, 0.3, "Deploy in 5 minutes  |  Zero infrastructure  |  AWS / Azure / GCP", size=13, color=GRAY, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 2: What It Does — 14 Tabs Overview
# =====================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2)
bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

add_text(s2, 0.6, 0.3, 12, 0.6, "What LHO Lite Covers", size=32, bold=True, color=WHITE)
add_text(s2, 0.6, 0.85, 12, 0.4, "14 interactive dashboard tabs providing complete workspace visibility", size=14, color=GRAY)

# 4 category cards
groups = [
    ("OVERVIEW", GREEN, [
        "Executive Summary — Security grade, cost, queries, findings at a glance",
        "Workspace Overview — Full inventory of users, clusters, catalogs, config flags",
    ]),
    ("SECURITY", RED, [
        "Compliance — HIPAA, FedRAMP, SOC 2, RBAC assessments with NIST mapping",
        "Architecture — Mermaid topology diagrams of workspace components",
    ]),
    ("OPERATIONS", BLUE, [
        "Infrastructure — Cluster & warehouse inventory with encryption/auto-stop status",
        "Spend Overview — 90-day cost trending by category, tag, and daily breakdown",
        "Workflows — Job run performance with success/failure rates and cost attribution",
        "Cost Explorer — Drill-down to billing line items by SKU, warehouse, job, notebook",
        "Cost Details — Cost estimation breakdown by user and compute category",
    ]),
    ("DATA & ACTIVITY", PURPLE, [
        "Apps & Models — Databricks Apps inventory and model serving endpoint pricing",
        "Table Inventory — Data asset discovery with storage analysis and search",
        "User Activity — Per-user query stats, data read volumes, compute time",
        "Daily Trends — Time-series analysis with heatmaps and dual-axis charts",
        "DBU Pricing — Current Databricks pricing reference by SKU category",
    ]),
]

y = 1.45
for title, color, items in groups:
    card_h = 0.25 + len(items) * 0.28
    add_card(s2, 0.5, y, 12.333, card_h)
    add_text(s2, 0.7, y + 0.05, 2, 0.3, title, size=11, bold=True, color=color)
    for j, item in enumerate(items):
        parts = item.split(" — ")
        tf = add_text(s2, 2.5, y + 0.05 + j * 0.28, 9.5, 0.28, "", size=11, color=WHITE)
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = parts[0]
        run1.font.bold = True
        run1.font.size = Pt(11)
        run1.font.color.rgb = WHITE
        run1.font.name = "Calibri"
        if len(parts) > 1:
            run2 = p.add_run()
            run2.text = f" — {parts[1]}"
            run2.font.size = Pt(11)
            run2.font.color.rgb = GRAY
            run2.font.name = "Calibri"
    y += card_h + 0.12

footer(s2)

# =====================================================================
# SLIDE 3: Data Sources & Critical Factors
# =====================================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3)
bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

add_text(s3, 0.6, 0.3, 12, 0.6, "Data Sources & Critical Security Factors", size=32, bold=True, color=WHITE)

# Left column: System tables
add_bullet_card(s3, 0.5, 1.2, 5.8, 2.6, "SYSTEM TABLES (Real Billing & Query Data)", [
    "system.billing.usage — 90-day cost by product, tag, job, warehouse, notebook",
    "system.billing.list_prices — Current DBU pricing for all SKUs",
    "system.query.history — 30-day user queries with status, GB read, duration",
    "system.compute.warehouse_events — Warehouse scaling events and utilization",
    "system.lakeflow.job_run_timeline — Job runs with success/failure and duration",
    "system.information_schema.tables — Table catalog with storage sizes",
], title_color=BLUE)

# Right column: REST APIs
add_bullet_card(s3, 7.0, 1.2, 5.8, 2.6, "REST APIs (22+ Endpoints)", [
    "SCIM — Users, Groups, Service Principals (identity & admin analysis)",
    "Clusters & Warehouses — Compute inventory, encryption, auto-termination",
    "Unity Catalog — Catalogs, metastores, credentials, shares, locations",
    "Workspace Config — 13 security flags (export, download, tokens, IP lists)",
    "Token Management — PAT lifecycle, expiration monitoring",
    "Apps & Serving — Databricks Apps + model endpoint pricing",
], title_color=GREEN)

# Bottom: Security findings
add_text(s3, 0.6, 4.1, 12, 0.4, "16 Automated Security Findings", size=20, bold=True, color=RED)

findings = [
    ("CRITICAL", RED, [
        "Hardcoded credentials in init scripts (AWS keys, passwords, Azure secrets)",
        "Excessive admin accounts (>20% admin ratio)",
        "Non-GovCloud deployment for regulated workloads",
    ]),
    ("HIGH", YELLOW, [
        "No IP access lists (open network boundary)",
        "Token lifetime exceeding 90 days",
        "Cluster disk encryption disabled",
        "Legacy security modes (data_security_mode=NONE)",
        "Missing session management controls",
    ]),
    ("MEDIUM", BLUE, [
        "External/non-org email domains in workspace",
        "Personal email accounts (gmail, outlook)",
        "Insufficient secret scopes for credential management",
        "Delta Sharing with external recipients",
        "Clusters without auto-termination",
    ]),
]

x = 0.5
for sev, color, items in findings:
    w = 4.1
    add_card(s3, x, 4.55, w, 2.2)
    add_text(s3, x + 0.15, 4.6, w - 0.3, 0.3, sev, size=13, bold=True, color=color)
    tf = add_text(s3, x + 0.15, 4.9, w - 0.3, 1.8, "", size=10, color=WHITE)
    for item in items:
        add_para(tf, f"  {item}", size=10, color=GRAY, spacing_before=3)
    x += w + 0.15

footer(s3)

# =====================================================================
# SLIDE 4: Ease of Use & Deployment
# =====================================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4)
bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

add_text(s4, 0.6, 0.3, 12, 0.6, "Deploy in 5 Minutes, Zero Infrastructure", size=32, bold=True, color=WHITE)
add_text(s4, 0.6, 0.85, 12, 0.4, "One notebook. One click. Full workspace visibility.", size=14, color=GRAY)

# 4 step cards
steps = [
    ("1", "Import Notebook", "Import the installer notebook\ninto your Databricks workspace.\nNo Terraform, no CI/CD, no\nexternal servers required.", BLUE),
    ("2", "Run All", "Enter your license key and\nclick Run All. The installer\ndownloads code, creates the\napp, and grants permissions.", GREEN),
    ("3", "Confirm Setup", "Open the app URL. The admin\npage appears with pre-filled\nconfig. Click Save & Start\nto begin data collection.", YELLOW),
    ("4", "Dashboard Ready", "Data collection completes in\n2-4 minutes. Full 14-tab\ndashboard is live. Set auto-\nrefresh for continuous monitoring.", PURPLE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 3.2
    add_card(s4, x, 1.5, 2.9, 2.8)
    # Number circle
    circ = s4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.05), Inches(1.65), Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf_c = circ.text_frame
    tf_c.paragraphs[0].text = num
    tf_c.paragraphs[0].font.size = Pt(24)
    tf_c.paragraphs[0].font.bold = True
    tf_c.paragraphs[0].font.color.rgb = BG_DARK
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_c.word_wrap = False

    add_text(s4, x + 0.15, 2.4, 2.6, 0.35, title, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, x + 0.15, 2.8, 2.6, 1.2, desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Key features row
add_text(s4, 0.6, 4.6, 12, 0.4, "Key Platform Features", size=20, bold=True, color=WHITE)

features = [
    ("Universal Cloud Support", "Works on AWS, Azure, and GCP\nDatabricks workspaces with\nauto-detection", BLUE),
    ("License-Gated Access", "Enterprise licensing with remote\nvalidation, expiration, key rotation,\nand 48-hour grace period", GREEN),
    ("Scheduled Refresh", "Manual, hourly, daily, or weekly\ndata collection. Auto-refresh\nkeeps dashboards current.", YELLOW),
    ("Excel Export", "One-click security and usage\nreport downloads for stakeholders\nwho don't access the dashboard", PURPLE),
]

for i, (title, desc, color) in enumerate(features):
    x = 0.5 + i * 3.2
    add_card(s4, x, 5.1, 2.9, 1.7)
    add_text(s4, x + 0.15, 5.2, 2.6, 0.3, title, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s4, x + 0.15, 5.55, 2.6, 1.0, desc, size=10, color=GRAY, align=PP_ALIGN.CENTER)

footer(s4)

# =====================================================================
# SLIDE 5: Value & ROI
# =====================================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5)
bar = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

add_text(s5, 0.6, 0.3, 12, 0.6, "Business Value & ROI", size=32, bold=True, color=WHITE)

# Three value pillars
pillars = [
    ("Cost Optimization", "$30K - $250K+/yr", "12-25% reduction in Databricks spend\nthrough waste identification, right-sizing,\nand workload optimization", BLUE, [
        "Identify idle warehouses & clusters",
        "Per-job cost attribution to owners",
        "Untagged workload detection",
        "Daily cost trend anomaly spotting",
    ]),
    ("Security & Compliance", "Risk Reduction", "Continuous monitoring replaces\nquarterly audits and point-in-time\nconsulting engagements", RED, [
        "Credential exposure detection",
        "HIPAA / FedRAMP / SOC 2 / RBAC",
        "NIST-mapped security findings",
        "16-point automated assessment",
    ]),
    ("Operational Efficiency", "55-111 hrs/month saved", "Replace manual data gathering\nwith automated, always-current\ndashboards and exports", GREEN, [
        "One-click Excel reports",
        "Real-time cost & usage dashboards",
        "User activity & job monitoring",
        "Infrastructure health inventory",
    ]),
]

for i, (title, metric, desc, color, items) in enumerate(pillars):
    x = 0.5 + i * 4.2
    add_card(s5, x, 1.1, 3.9, 3.6)
    add_text(s5, x + 0.2, 1.2, 3.5, 0.3, title, size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s5, x + 0.2, 1.55, 3.5, 0.35, metric, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s5, x + 0.2, 1.95, 3.5, 0.9, desc, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    tf = add_text(s5, x + 0.2, 2.85, 3.5, 1.5, "", size=10, color=WHITE)
    for item in items:
        add_para(tf, f"  {item}", size=10, color=WHITE, spacing_before=3)

# ROI summary bar
add_card(s5, 0.5, 5.0, 12.333, 1.2, fill_color=RGBColor(0x10, 0x18, 0x25))

roi_items = [
    ("10-30x", "ROI within 90 days"),
    ("<30 days", "Payback period"),
    ("5 minutes", "Time to deploy"),
    ("0 hours", "Infrastructure setup"),
    ("$0", "Additional cloud cost"),
]

for i, (val, label) in enumerate(roi_items):
    x = 0.9 + i * 2.45
    add_text(s5, x, 5.1, 2.2, 0.5, val, size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s5, x, 5.55, 2.2, 0.3, label, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Bottom tagline
add_text(s5, 0.5, 6.4, 12.333, 0.4,
         "LHO Lite: Deploy once, monitor continuously, optimize always.",
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(s5)

# Save
output_path = "/Users/darkstar33/Documents/testing2/lho-lite/docs/LHO_Lite_Overview_Deck.pptx"
prs.save(output_path)
print(f"Deck saved to {output_path}")
